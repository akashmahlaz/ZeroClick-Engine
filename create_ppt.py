from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# ─── Color Palette ──────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x0D, 0x2B, 0x55)   # #0D2B55
MID_BLUE    = RGBColor(0x1A, 0x5F, 0xA4)   # #1A5FA4
LIGHT_BLUE  = RGBColor(0x5B, 0xAE, 0xE8)   # #5BAEE8
ACCENT      = RGBColor(0xF5, 0xA6, 0x23)   # #F5A623
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY   = RGBColor(0x2D, 0x2D, 0x2D)
LIGHT_GRAY  = RGBColor(0xF4, 0xF6, 0xF9)
GREEN       = RGBColor(0x27, 0xAE, 0x60)

# ─── Helper Functions ──────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox

def add_bullet_text(slide, items, left, top, width, height,
                   font_size=16, color=DARK_GRAY, bullet_color=MID_BLUE):
    from pptx.util import Pt as Pt2
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(6)
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = f"• {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color

def slide_header(slide, title, subtitle=None):
    # Top bar
    add_rect(slide, 0, 0, 13.33, 1.1, DARK_BLUE)
    add_text_box(slide, title, 0.5, 0.2, 12, 0.7,
                 font_size=32, bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, subtitle, 0.5, 0.75, 12, 0.35,
                     font_size=14, bold=False, color=LIGHT_BLUE)
    # Accent bar
    add_rect(slide, 0, 1.1, 13.33, 0.06, ACCENT)

def content_card(slide, left, top, width, height, title, items, title_color=MID_BLUE, bg_color=WHITE):
    from pptx.util import Pt as Pt2
    # Card background
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(0xE0, 0xE8, 0xF0)
    shape.line.width = Pt(1)

    # Title bar
    title_bar = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(0.5))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = title_color
    title_bar.line.fill.background()

    add_text_box(slide, title, left+0.15, top+0.07, width-0.3, 0.38,
                 font_size=15, bold=True, color=WHITE)

    # Items
    txBox = slide.shapes.add_textbox(Inches(left+0.2), Inches(top+0.55),
                                      Inches(width-0.4), Inches(height-0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = f"→ {item}"
        run.font.size = Pt(13)
        run.font.color.rgb = DARK_GRAY

# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])

# Background
add_rect(slide1, 0, 0, 13.33, 7.5, DARK_BLUE)
# Accent line
add_rect(slide1, 0, 5.5, 13.33, 0.06, ACCENT)

add_text_box(slide1, "RxNetwork × Adverge", 1, 1.5, 11, 0.8,
             font_size=18, bold=False, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
add_text_box(slide1, "In-Content Text-Only Unit", 0.5, 2.3, 12, 1.2,
             font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide1, "Technical Delivery Proposal", 0.5, 3.5, 12, 0.6,
             font_size=24, bold=False, color=ACCENT, align=PP_ALIGN.CENTER)

add_text_box(slide1, "Prepared by: [Your Company Name]", 1, 5.7, 11, 0.5,
             font_size=14, bold=False, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
add_text_box(slide1, "April 2026", 1, 6.3, 11, 0.5,
             font_size=13, bold=False, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ═══════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide2, "The Problem", "Why RxNetwork needs a new revenue strategy")

# Left card
content_card(slide2, 0.4, 1.5, 6, 2.8, "Zero-Click AI Search",
             ["LLMs (ChatGPT, Gemini, Copilot) answer directly on search page",
              "Users no longer click through to publisher websites",
              "RxNetwork sites losing referral traffic measurable declines",
              "Ad inventory and revenue dropping as pageviews fall"])

content_card(slide2, 0.4, 4.5, 6, 2.5, "The Opportunity",
             ["Brands want visibility inside AI-generated answers",
              "Healthcare advertisers need verified HCP audiences",
              "Text-only format is LLM-crawlable and compliant",
              "Premium inventory with guaranteed targeting"], title_color=GREEN)

# Right side — big callout
add_rect(slide2, 6.7, 1.5, 6.2, 5.5, RGBColor(0x1A, 0x5F, 0xA4))
add_text_box(slide2, "\"LLM crawlers can legitimately capture these snippets and surface them—along with brand mentions—in zero-click search summaries.\"",
             7, 2.2, 5.5, 2, font_size=17, bold=False, color=WHITE, align=PP_ALIGN.LEFT)
add_text_box(slide2, "The solution: place factual, medically accurate brand mentions directly inside article content — content that AI can safely ingest and surface.",
             7, 4.5, 5.5, 1.5, font_size=14, bold=False, color=LIGHT_BLUE, align=PP_ALIGN.LEFT)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — OUR SOLUTION (HIGH LEVEL)
# ═══════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide3, "Our Solution", "A complete in-content text unit delivery system")

# Three cards side by side
content_card(slide3, 0.4, 1.5, 4, 2.5, "JavaScript Tag",
             ["Lightweight tag (under 5KB)",
              "Paste once on every article page",
              "Loads asynchronously — no page lag",
              "Matches page CSS automatically"], title_color=DARK_BLUE)

content_card(slide3, 4.65, 1.5, 4, 2.5, "Campaign Engine",
             ["Matches NPI specialty + page category",
              "Serves only verified HCPs",
              "Dynamically selects right text snippet",
              "No fallback to non-HCP traffic"], title_color=MID_BLUE)

content_card(slide3, 8.9, 1.5, 4, 2.5, "Reporting Dashboard",
             ["Impressions, geo, device breakdown",
              "NPI specialty delivery reports",
              "Campaign performance at a glance",
              "CSV export for buyers"], title_color=GREEN)

# Bottom explanation
add_rect(slide3, 0.4, 4.3, 12.5, 2.8, LIGHT_GRAY)
add_text_box(slide3, "How it works — end to end:", 0.7, 4.5, 12, 0.4,
             font_size=15, bold=True, color=DARK_BLUE)
steps = [
    "Doctor visits article page",
    "JS tag fires → checks NPI",
    "Server matches campaign",
    "Text snippet injected in-content",
    "Impression logged + reported"
]
for i, step in enumerate(steps):
    x = 0.7 + i * 2.5
    add_rect(slide3, x, 5.1, 2.2, 0.65, MID_BLUE)
    add_text_box(slide3, f"Step {i+1}\n{step}", x+0.1, 5.12, 2, 0.6,
                 font_size=11, bold=False, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 4:
        add_text_box(slide3, "→", x+2.15, 5.15, 0.4, 0.5,
                     font_size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — SYSTEM ARCHITECTURE
# ═══════════════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide4, "System Architecture", "All components and how they connect")

boxes = [
    ("RxNetwork Sites", "Medical content sites\nJS tag installed on\narticle pages", 0.4, 1.4, MID_BLUE),
    ("Client Browser", "Visitor's browser\nloads JS tag\nchecks NPI cookie", 0.4, 4.2, MID_BLUE),
    ("Your Backend\n(FastAPI + Python)", "Campaign matching\nNPI verification\nSnippet delivery\nImpression logging", 4.1, 1.4, DARK_BLUE),
    ("Redis Cache", "Fast lookup\nReduce DB load\nSession data", 7.8, 4.2, GREEN),
    ("Admin Dashboard\n(React + Next.js)", "Create campaigns\nView reports\nManage targeting", 11.2, 1.4, ACCENT),
    ("NPI Identity\nProviders", "Doceree / DeepIntent\nor RxNetwork\npass-through", 11.2, 4.2, ACCENT),
]

for title, desc, left, top, color in boxes:
    add_rect(slide4, left, top, 3, 2.5, color)
    add_text_box(slide4, title, left+0.1, top+0.15, 2.8, 0.6,
                 font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide4, desc, left+0.1, top+0.8, 2.8, 1.5,
                 font_size=12, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

# Arrows (text-based arrows)
arrows = [
    (3.45, 2.4, 0.65, 0),  # Site → Backend
    (7.55, 2.4, 0.55, 0),  # Backend → DB
    (10.95, 2.4, 0.25, 0), # Backend → Dashboard
    (3.45, 5.0, 0.65, 0),  # Browser → Backend
    (10.95, 5.0, 0.25, 0), # Browser → NPI
    (7.55, 5.0, 0.25, 0),  # DB → Cache
]
for ax, ay, aw, ah in arrows:
    add_rect(slide4, ax, ay, aw, 0.08, ACCENT)

add_text_box(slide4, "HTTPS all connections | Redis caching | PostgreSQL storage | Async JS loading",
             0.4, 7.0, 12.5, 0.4, font_size=12, bold=False, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — TECH STACK
# ═══════════════════════════════════════════════════════════════
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide5, "Technology Stack", "Purpose + Why We Chose Each Technology")

rows = [
    ("JavaScript (Vanilla)", "Runs on all RxNetwork sites", "Lightweight, no framework, works everywhere, async load"),
    ("Python (FastAPI)", "Backend API + campaign matching", "Fast, async, handles many concurrent requests, great for ad-tech"),
    ("PostgreSQL", "Primary database", "Relational, reliable, handles complex queries for reporting"),
    ("Redis", "Caching layer", "Fast lookups, reduces DB load, handles session data"),
    ("React + Next.js", "Admin dashboard", "Modern, fast, great UI components, easy to deploy on Cloudflare"),
    ("Cloudflare Workers", "Hosting (optional)", "Edge deployment, fast globally, scales automatically, cost-effective"),
    ("NPI Registry API / Third-party", "Doctor identity verification", "Either CMS NPI registry or providers like Doceree, DeepIntent"),
]

add_rect(slide5, 0.4, 1.35, 12.5, 0.5, RGBColor(0xE8, 0xF0, 0xF8))
headers = ["Technology", "Purpose", "Why This Choice"]
for i, h in enumerate(headers):
    add_text_box(slide5, h, 0.5 + i*4.2, 1.4, 4, 0.4,
                 font_size=13, bold=True, color=DARK_BLUE)

for row_i, (tech, purpose, why) in enumerate(rows):
    y = 1.95 + row_i * 0.72
    bg = WHITE if row_i % 2 == 0 else RGBColor(0xF4, 0xF6, 0xF9)
    add_rect(slide5, 0.4, y, 12.5, 0.68, bg)
    add_text_box(slide5, tech, 0.5, y+0.12, 4, 0.45, font_size=13, bold=True, color=MID_BLUE)
    add_text_box(slide5, purpose, 4.5, y+0.12, 4, 0.45, font_size=13, bold=False, color=DARK_GRAY)
    add_text_box(slide5, why, 8.7, y+0.12, 4, 0.45, font_size=13, bold=False, color=DARK_GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — HOW IT WORKS (WORKFLOW)
# ═══════════════════════════════════════════════════════════════
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide6, "How the System Works", "Step-by-step flow from page load to reporting")

steps_data = [
    ("1", "User Loads Article Page", MID_BLUE,
     "Doctor visits RxNetwork site → JS tag fires automatically in background"),
    ("2", "NPI Eligibility Check", DARK_BLUE,
     "Tag checks NPI status via Adverge identity graph or RxNetwork pass-through"),
    ("3", "NO NPI → No Op", RGBColor(0xC0, 0x39, 0x2B),
     "If NPI missing/invalid → snippet not requested, nothing shown"),
    ("4", "YES NPI → Request Snippet", GREEN,
     "Tag calls your backend with: site ID, page URL, unit category, NPI specialty"),
    ("5", "Campaign Matching", RGBColor(0xF5, 0xA6, 0x23),
     "Backend matches: page category + NPI specialty + active campaign budget"),
    ("6", "HTML Snippet Delivered", MID_BLUE,
     "1–10 lines of text returned in HTML format, styled to match host page CSS"),
    ("7", "Snippet Injected + Reported", DARK_BLUE,
     "JS inserts snippet into article → impression logged → dashboard updated"),
]

for i, (num, title, color, desc) in enumerate(steps_data):
    row = i // 2
    col = i % 2
    x = 0.4 + col * 6.5
    y = 1.4 + row * 1.85

    add_rect(slide6, x, y, 6, 1.7, color)
    add_text_box(slide6, f"Step {num}", x+0.15, y+0.1, 5.5, 0.3,
                 font_size=11, bold=False, color=LIGHT_BLUE)
    add_text_box(slide6, title, x+0.15, y+0.4, 5.5, 0.45,
                 font_size=15, bold=True, color=WHITE)
    add_text_box(slide6, desc, x+0.15, y+0.95, 5.5, 0.6,
                 font_size=12, bold=False, color=RGBColor(0xE0, 0xE8, 0xF0))

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — TARGETING DETAILS
# ═══════════════════════════════════════════════════════════════
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide7, "NPI-Based Targeting + Medical Categories", "How campaigns are matched and delivered")

# Left: NPI Targeting
add_rect(slide7, 0.4, 1.4, 6, 5.5, LIGHT_GRAY)
add_rect(slide7, 0.4, 1.4, 6, 0.6, DARK_BLUE)
add_text_box(slide7, "NPI Targeting Logic", 0.55, 1.47, 5.7, 0.45,
             font_size=16, bold=True, color=WHITE)

npi_items = [
    "Every US doctor has a unique 10-digit NPI number",
    "Only NPI-verified HCPs receive the text snippet",
    "Non-verified visitors see no snippet (unit collapses)",
    "NPI data passed as hashed/pseudonymous identifier",
    "No PHI (Protected Health Information) stored or transmitted",
    "Integration methods: login-based, third-party cookie, or NPI pass-through",
    "Specialty extracted from NPI for targeting segmentation",
    "Frequency capping supported for future enhancement",
]
add_bullet_text(slide7, npi_items, 0.55, 2.15, 5.7, 4.5, font_size=13, color=DARK_GRAY)

# Right: Medical Categories
add_rect(slide7, 6.8, 1.4, 6.1, 5.5, LIGHT_GRAY)
add_rect(slide7, 6.8, 1.4, 6.1, 0.6, MID_BLUE)
add_text_box(slide7, "Medical Category Assignment", 6.95, 1.47, 5.8, 0.45,
             font_size=16, bold=True, color=WHITE)

cat_items = [
    "Each page placement is assigned one or more categories",
    "Campaigns targeted by: NPI specialty + category match",
    "Only campaigns matching BOTH specialty AND category are eligible",
    "No fallback to non-HCP or non-matching campaigns",
    "Supported categories:",
    "  — Pediatrics, OBGYN, Geriatrics",
    "  — Pain, Rheumatology, Gastroenterology",
    "  — Additional specialties as needed",
]
add_bullet_text(slide7, cat_items, 6.95, 2.15, 5.7, 4.5, font_size=13, color=DARK_GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — CLOUDFLARE WORKERS SECTION
# ═══════════════════════════════════════════════════════════════
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide8, "Cloudflare Workers Deployment", "Why we recommend it and how it fits the architecture")

# Left explanation
add_rect(slide8, 0.4, 1.4, 7.5, 5.5, LIGHT_GRAY)
add_rect(slide8, 0.4, 1.4, 7.5, 0.6, RGBColor(0xF5, 0xA6, 0x23))
add_text_box(slide8, "What is Cloudflare Workers?", 0.55, 1.47, 7.2, 0.45,
             font_size=16, bold=True, color=WHITE)

cf_items = [
    "Cloudflare Workers runs your code on edge servers worldwide",
    "Instead of one server, your backend runs in 300+ data centers",
    "When a doctor in New York visits RxNetwork → code runs locally (fast)",
    "When a doctor in London visits → same, code runs locally there",
    "No single point of failure — if one server fails, another takes over",
    "Automatically scales — handles 100 visits or 10 million visits equally",
    "No server management — Cloudflare handles infrastructure completely",
    "Free tier available; paid plans start at $5/month for high traffic",
]
add_bullet_text(slide8, cf_items, 0.55, 2.15, 7.2, 4.5, font_size=13, color=DARK_GRAY)

# Right: Comparison
add_rect(slide8, 8.2, 1.4, 4.7, 2.4, MID_BLUE)
add_text_box(slide8, "Traditional Server", 8.35, 1.5, 4.4, 0.35,
             font_size=14, bold=True, color=WHITE)
add_text_box(slide8, "One physical location\nSingle point of failure\nManual scaling needed\nHigher latency for global users\nServer maintenance required",
             8.35, 1.95, 4.4, 1.7, font_size=12, bold=False, color=RGBColor(0xE0, 0xE8, 0xF0))

add_rect(slide8, 8.2, 4.0, 4.7, 2.9, GREEN)
add_text_box(slide8, "Cloudflare Workers ✓", 8.35, 4.1, 4.4, 0.35,
             font_size=14, bold=True, color=WHITE)
add_text_box(slide8, "300+ edge locations worldwide\nGlobally distributed, no single failure\nAuto-scaling, no manual intervention\nLow latency for all global users\nNo server maintenance needed",
             8.35, 4.55, 4.4, 2.2, font_size=12, bold=False, color=WHITE)

# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — CLOUDFLARE ARCHITECTURE DETAIL
# ═══════════════════════════════════════════════════════════════
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide9, "Cloudflare Workers — Architecture Details", "How each component deploys on Cloudflare")

cloud_components = [
    ("Workers (FastAPI backend)", "Python/FastAPI code runs here\nHandles all API calls\nSnippet matching logic\nCampaign routing\nImpression logging", DARK_BLUE, 0.4, 1.4),
    ("Cloudflare KV", "Key-value storage\nStores campaign data\nSession caching\nFast global access", RGBColor(0x27, 0xAE, 0x60), 0.4, 4.4),
    ("Cloudflare D1", "SQLite-based database\nStores impressions, reports\nSQL queries for dashboard\nGlobally replicated", MID_BLUE, 7.3, 1.4),
    ("R2 Object Storage", "Stores text snippet content\nFallback for large campaigns\nCDN-backed delivery", ACCENT, 7.3, 4.4),
]

for title, desc, color, left, top in cloud_components:
    add_rect(slide9, left, top, 5.6, 2.7, color)
    add_text_box(slide9, title, left+0.15, top+0.15, 5.3, 0.5,
                 font_size=14, bold=True, color=WHITE)
    add_text_box(slide9, desc, left+0.15, top+0.75, 5.3, 1.8,
                 font_size=12, bold=False, color=WHITE)

add_text_box(slide9, "All connections use HTTPS | Workers run at edge | D1/KV replicate globally automatically",
             0.4, 7.15, 12.5, 0.35, font_size=12, bold=False, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — PHASES & TIMELINE
# ═══════════════════════════════════════════════════════════════
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide10, "Delivery Phases & Timeline", "Phased approach for risk-free validation and scaling")

phases = [
    ("Phase 1 — MVP", "8–10 weeks", DARK_BLUE,
     ["Core backend (FastAPI + database)",
      "JavaScript tag (single site)",
      "Basic NPI integration",
      "Campaign matching logic",
      "Single-site pilot testing"]),
    ("Phase 2 — Rollout", "6–8 weeks", MID_BLUE,
     ["Multi-site deployment",
      "Full category targeting",
      "Advanced NPI segmentation",
      "Frequency capping",
      "Performance optimization"]),
    ("Phase 3 — Dashboard", "6–8 weeks", GREEN,
     ["React/Next.js admin dashboard",
      "Reporting + CSV export",
      "Campaign management UI",
      "Security hardening (HTTPS, hashing)",
      "Client training + handover"]),
]

for i, (title, weeks, color, items) in enumerate(phases):
    x = 0.4 + i * 4.35
    add_rect(slide10, x, 1.4, 4.1, 0.7, color)
    add_text_box(slide10, title, x+0.15, 1.47, 3.8, 0.35,
                 font_size=15, bold=True, color=WHITE)
    add_text_box(slide10, f"{weeks}", x+0.15, 1.92, 3.8, 0.35,
                 font_size=13, bold=False, color=ACCENT)
    add_rect(slide10, x, 2.15, 4.1, 4.8, LIGHT_GRAY)
    add_bullet_text(slide10, items, x+0.2, 2.3, 3.7, 4.5, font_size=13, color=DARK_GRAY)

add_text_box(slide10, "Total Timeline: 20–26 weeks  |  Phase 1 delivers working system  |  Subsequent phases build on validated foundation",
             0.4, 7.1, 12.5, 0.4, font_size=13, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — INVESTMENT
# ═══════════════════════════════════════════════════════════════
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide11, "Investment", "Transparent, phase-based pricing designed for fast validation")

# Main pricing box
add_rect(slide11, 0.4, 1.4, 12.5, 2.2, DARK_BLUE)
add_text_box(slide11, "Total Estimated Investment: $18,000 – $28,000", 0.7, 1.6, 11.9, 0.7,
             font_size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide11, "Phase-based delivery — pay as we deliver. No large upfront commitment.",
             0.7, 2.35, 11.9, 0.5, font_size=16, bold=False, color=ACCENT, align=PP_ALIGN.CENTER)

phase_pricing = [
    ("Phase 1 — MVP", "$8,000 – $12,000", "Core engine, JS tag, single-site pilot"),
    ("Phase 2 — Rollout", "$5,000 – $8,000", "Multi-site, advanced targeting, optimization"),
    ("Phase 3 — Dashboard", "$5,000 – $8,000", "Reporting UI, compliance, training"),
]

for i, (title, price, desc) in enumerate(phase_pricing):
    x = 0.4 + i * 4.35
    add_rect(slide11, x, 3.8, 4.1, 2.8, MID_BLUE)
    add_text_box(slide11, title, x+0.15, 3.95, 3.8, 0.45,
                 font_size=15, bold=True, color=WHITE)
    add_text_box(slide11, price, x+0.15, 4.5, 3.8, 0.55,
                 font_size=22, bold=True, color=ACCENT)
    add_text_box(slide11, desc, x+0.15, 5.2, 3.8, 1.2,
                 font_size=12, bold=False, color=RGBColor(0xE0, 0xE8, 0xF0))

add_text_box(slide11, "Market context: US agencies charge $80,000–$130,000 for equivalent systems. Our lean team delivers the same outcome at a fraction of the cost.",
             0.4, 6.75, 12.5, 0.6, font_size=12, bold=False, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — WHY CHOOSE US
# ═══════════════════════════════════════════════════════════════
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide12, "Why Work With Us", "Smart positioning: not a freelancer, not an expensive agency")

comparisons = [
    ("Freelancer", "$3K–$8K", "Limited scope, no team, single point of failure, hard to scale", RGBColor(0xC0, 0x39, 0x2B)),
    ("US Agency", "$80K–$130K", "Full scope, but expensive, slow, and overkill for this stage", RGBColor(0x8E, 0x44, 0xAD)),
    ("✅ Us (Lean Team)", "$18K–$28K", "Right-sized team, full system, faster delivery, cost-effective", GREEN),
]

for i, (title, price, risk, color) in enumerate(comparisons):
    x = 0.4 + i * 4.35
    add_rect(slide12, x, 1.5, 4.1, 4.0, color)
    add_text_box(slide12, title, x+0.15, 1.65, 3.8, 0.5,
                 font_size=16, bold=True, color=WHITE)
    add_text_box(slide12, price, x+0.15, 2.25, 3.8, 0.55,
                 font_size=20, bold=True, color=ACCENT)
    add_text_box(slide12, risk, x+0.15, 3.0, 3.8, 2.2,
                 font_size=13, bold=False, color=WHITE)

benefits = [
    "Full AdTech system — not just a script",
    "Phased delivery — validate before full spend",
    "Healthcare + ad-tech domain knowledge",
    "Cloud-native, globally distributed architecture",
    "Transparent pricing — no hidden costs",
]
add_rect(slide12, 0.4, 5.7, 12.5, 1.6, LIGHT_GRAY)
add_text_box(slide12, "What you get:", 0.55, 5.8, 12, 0.35,
             font_size=13, bold=True, color=DARK_BLUE)
add_bullet_text(slide12, benefits, 0.55, 6.2, 12, 1.0, font_size=13, color=DARK_GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — NEXT STEPS
# ═══════════════════════════════════════════════════════════════
slide13 = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide13, "Next Steps", "Let's get started — here's what happens after you approve")

steps = [
    ("Week 1", "Discovery Call", "Clarify NPI integration method, ad server stack, and site inventory"),
    ("Week 2", "Contract + MSA", "Sign agreement, collect 30% upfront payment"),
    ("Week 3–4", "Design Phase", "Database schema, API design, architecture documentation"),
    ("Week 5–10", "Phase 1 Build", "Backend, JS tag, single-site pilot live"),
    ("Week 11–12", "Phase 1 Review", "Validate with client, collect feedback, prepare Phase 2"),
]

for i, (week, title, desc) in enumerate(steps):
    y = 1.5 + i * 1.1
    add_rect(slide13, 0.4, y, 2.2, 0.9, MID_BLUE)
    add_text_box(slide13, week, 0.5, y+0.1, 2, 0.35, font_size=13, bold=True, color=WHITE)
    add_text_box(slide13, title, 0.5, y+0.5, 2, 0.35, font_size=12, bold=False, color=LIGHT_BLUE)
    add_rect(slide13, 2.8, y, 10.1, 0.9, LIGHT_GRAY)
    add_text_box(slide13, desc, 3.0, y+0.25, 9.7, 0.5, font_size=13, bold=False, color=DARK_GRAY)
    if i < 4:
        add_text_box(slide13, "▼", 1.2, y+0.9, 0.5, 0.3,
                     font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

add_rect(slide13, 0.4, 6.95, 12.5, 0.4, GREEN)
add_text_box(slide13, "Ready to start? Let's schedule a technical discovery call this week. 👇",
             0.7, 6.97, 12, 0.35, font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════
prs.save("RxNetwork_Proposal.pptx")
print("✅ PowerPoint saved as RxNetwork_Proposal.pptx")