"""Generate RxNetwork In-Content Unit Proposal — Clean White Professional Design"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Professional White Color Palette ──
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)
BODY_TEXT = RGBColor(0x37, 0x41, 0x51)
MUTED = RGBColor(0x6B, 0x72, 0x80)
CARD_BG = RGBColor(0xF1, 0xF5, 0xF9)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)
PRIMARY = RGBColor(0x25, 0x63, 0xEB)
PRIMARY_LIGHT = RGBColor(0xDB, 0xEA, 0xFE)
SUCCESS = RGBColor(0x05, 0x96, 0x69)
SUCCESS_LIGHT = RGBColor(0xD1, 0xFA, 0xE5)
WARNING = RGBColor(0xD9, 0x77, 0x06)
WARNING_LIGHT = RGBColor(0xFE, 0xF3, 0xC7)
DANGER = RGBColor(0xDC, 0x26, 0x26)
DANGER_LIGHT = RGBColor(0xFE, 0xE2, 0xE2)
ACCENT = RGBColor(0x7C, 0x3A, 0xED)
ACCENT_LIGHT = RGBColor(0xED, 0xE9, 0xFE)

TOTAL_SLIDES = 11

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, font_size=18,
             color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline(slide, left, top, width, height, lines, font_size=14,
                  color=BODY_TEXT, bold=False, alignment=PP_ALIGN.LEFT, line_spacing=1.2):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Calibri"
        p.space_after = Pt(font_size * (line_spacing - 1) + 2)
        p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=14, color=BODY_TEXT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  \u2022  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
    return txBox


def add_card(slide, left, top, width, height, fill_color=CARD_BG, border_color=BORDER):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def add_line(slide, left, top, width, color=PRIMARY, thickness=3):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(thickness / 72))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_divider(slide, left, top, width):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(0.01))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BORDER
    shape.line.fill.background()


def slide_header(slide, title, subtitle=None):
    add_bg(slide)
    add_line(slide, 1, 0.6, 1.2, PRIMARY, 4)
    add_text(slide, 1, 0.75, 10, 0.7, title, font_size=32, bold=True, color=DARK_TEXT)
    if subtitle:
        add_text(slide, 1, 1.25, 10, 0.4, subtitle, font_size=16, color=MUTED)
    add_divider(slide, 1, 1.65 if subtitle else 1.35, 11.3)


def page_number(slide, num):
    add_text(slide, 11.5, 7.1, 1.5, 0.3, f"{num} / {TOTAL_SLIDES}",
             font_size=10, color=MUTED, alignment=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════
# SLIDE 1: TITLE
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

top_bar = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.08))
top_bar.fill.solid()
top_bar.fill.fore_color.rgb = PRIMARY
top_bar.line.fill.background()

add_line(slide, 1.5, 2.6, 1.5, PRIMARY, 4)
add_text(slide, 1.5, 2.85, 10, 1, "In-Content Text-Only Unit",
         font_size=48, bold=True, color=DARK_TEXT)
add_text(slide, 1.5, 3.85, 10, 0.6,
         "Programmatic Delivery System for Healthcare Advertisers",
         font_size=24, color=PRIMARY)
add_text(slide, 1.5, 4.55, 10, 0.5,
         "Technical Proposal & Delivery Plan",
         font_size=18, color=MUTED)
add_divider(slide, 1.5, 5.3, 4)
add_text(slide, 1.5, 5.5, 5, 0.3,
         "Prepared for RxNetwork & Adverge",
         font_size=14, color=BODY_TEXT)
add_text(slide, 1.5, 5.85, 5, 0.3,
         "Prepared by Logician  \u2022  April 2026  \u2022  Confidential",
         font_size=12, color=MUTED)

# ════════════════════════════════════════════
# SLIDE 2: UNDERSTANDING YOUR CHALLENGE
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Understanding Your Challenge",
             "The shift to zero-click AI search is fundamentally changing healthcare publishing")
page_number(slide, 2)

add_card(slide, 1, 1.85, 5.5, 2.5, DANGER_LIGHT, DANGER)
add_text(slide, 1.3, 2.0, 5, 0.35, "The Problem",
         font_size=18, bold=True, color=DANGER)
add_bullet_list(slide, 1.3, 2.45, 5, 1.8, [
    "LLMs (Gemini, ChatGPT, Copilot) deliver answers directly on search pages",
    "Healthcare professionals no longer click through to publisher sites",
    "RxNetwork sites experiencing measurable traffic declines",
    "Fewer pageviews means shrinking ad inventory and revenue",
], font_size=13, color=DARK_TEXT)

add_card(slide, 6.8, 1.85, 5.5, 2.5, SUCCESS_LIGHT, SUCCESS)
add_text(slide, 7.1, 2.0, 5, 0.35, "The Opportunity",
         font_size=18, bold=True, color=SUCCESS)
add_bullet_list(slide, 7.1, 2.45, 5, 1.8, [
    "Pharma brands need visibility inside AI-summarized content",
    "Text-only format is inherently LLM-crawlable and compliant",
    "NPI-verified HCP audiences command premium CPMs",
    "Category-level targeting aligns with existing pharma budgets",
], font_size=13, color=DARK_TEXT)

add_card(slide, 1, 4.65, 11.3, 2.5)
add_text(slide, 1.3, 4.8, 10, 0.35, "Your Vision (As We Understand It)",
         font_size=18, bold=True, color=PRIMARY)
add_multiline(slide, 1.3, 5.25, 10.5, 1.6, [
    "Place 1\u201310 lines of medically accurate, non-promotional text with factual brand mentions directly inside article content.",
    "Serve these units exclusively to NPI-verified healthcare professionals visiting RxNetwork sites.",
    "Because the content is factual HTML integrated into the editorial flow, LLM crawlers can capture these snippets and surface them\u2014along with brand mentions\u2014in zero-click search summaries.",
    "Healthcare advertisers purchase targeted campaigns routed by medical category and NPI specialty.",
], font_size=13, color=BODY_TEXT)

# ════════════════════════════════════════════
# SLIDE 3: OUR APPROACH
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Our Approach",
             "A purpose-built delivery system \u2014 not a repurposed ad server")
page_number(slide, 3)

pillars = [
    ("Lightweight Integration",
     "A single JavaScript tag placed on article pages. Loads asynchronously, inherits host page CSS, and causes zero layout disruption. RxNetwork installs once \u2014 we handle the rest.",
     PRIMARY, PRIMARY_LIGHT),
    ("Precision Targeting",
     "Every impression is verified against NPI status. Campaigns are matched by medical specialty AND page category. No fallback to non-HCP traffic \u2014 every served unit reaches a verified healthcare professional.",
     SUCCESS, SUCCESS_LIGHT),
    ("Transparent Reporting",
     "Full programmatic reporting: impressions by site, campaign, specialty, geo, and device. Dashboard access for real-time monitoring plus CSV export for buyers and internal analysis.",
     WARNING, WARNING_LIGHT),
]

for i, (title, desc, color, bg) in enumerate(pillars):
    x = 0.7 + i * 4.1
    add_card(slide, x, 1.85, 3.85, 4.0, bg, color)
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x + 1.55), Inches(2.05), Inches(0.55), Inches(0.55))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = str(i + 1)
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_text(slide, x + 0.2, 2.75, 3.45, 0.4, title,
             font_size=16, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + 0.25, 3.25, 3.35, 2.4, desc,
             font_size=12, color=BODY_TEXT, alignment=PP_ALIGN.CENTER)

add_card(slide, 0.7, 6.15, 11.9, 0.8)
add_text(slide, 1, 6.25, 11, 0.5,
         "Key principle: The text snippet is pure HTML \u2014 not an image, not an iframe. "
         "This ensures search engines and LLM crawlers can read and index the content naturally, "
         "while NPI gating ensures only verified HCPs receive the unit on-page.",
         font_size=12, color=BODY_TEXT, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# SLIDE 4: WHAT WE WILL DELIVER
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "What We Will Deliver",
             "Four concrete deliverables that form a complete system")
page_number(slide, 4)

deliverables = [
    ("JavaScript Tag",
     "Lightweight async script (<5KB)\nthat RxNetwork embeds on article\npages. Handles NPI check, campaign\nrequest, and snippet injection\nwith zero page-load impact.",
     PRIMARY, PRIMARY_LIGHT),
    ("Campaign Matching Engine",
     "Backend server that receives\nrequests, verifies NPI status,\nmatches active campaigns by\nmedical category + specialty,\nand returns the right snippet.",
     SUCCESS, SUCCESS_LIGHT),
    ("Reporting Dashboard",
     "Web-based admin panel for\ncampaign management, performance\nreporting, audience delivery\nmetrics, and CSV export\nfor buyers and agencies.",
     WARNING, WARNING_LIGHT),
    ("Documentation & Training",
     "Complete deployment guide,\nAPI documentation, campaign\ncreation walkthrough, and\nlive training session for\nRxNetwork and Adverge teams.",
     ACCENT, ACCENT_LIGHT),
]

for i, (title, desc, color, bg) in enumerate(deliverables):
    x = 0.7 + i * 3.1
    add_card(slide, x, 1.85, 2.9, 3.8, bg, color)
    add_text(slide, x + 0.15, 2.05, 2.6, 0.35, title,
             font_size=15, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_line(slide, x + 0.5, 2.45, 1.9, color, 2)
    add_text(slide, x + 0.2, 2.65, 2.5, 3, desc,
             font_size=12, color=BODY_TEXT, alignment=PP_ALIGN.CENTER)

add_card(slide, 0.7, 5.95, 11.9, 1.2)
add_text(slide, 1, 6.05, 11, 0.35, "Reporting Includes",
         font_size=14, bold=True, color=PRIMARY)
add_text(slide, 1, 6.4, 11, 0.5,
         "Impressions served  \u2022  Viewability metrics  \u2022  Geo distribution  \u2022  "
         "Device breakdown  \u2022  NPI specialty segment  \u2022  Campaign ID  \u2022  "
         "Snippet ID  \u2022  Page URL / Site  \u2022  Downloadable CSV",
         font_size=12, color=BODY_TEXT)

# ════════════════════════════════════════════
# SLIDE 5: HOW IT WORKS
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "How It Works",
             "End-to-end flow from page load to impression reporting")
page_number(slide, 5)

steps = [
    ("1", "Page Load", "HCP visits an\narticle page on\nRxNetwork site", PRIMARY),
    ("2", "Tag Fires", "Lightweight JS\ntag loads async\nin background", PRIMARY),
    ("3", "NPI Check", "System verifies\nvisitor NPI status\nvia identity layer", SUCCESS),
    ("4", "Campaign\nMatch", "Engine matches\nspecialty + category\n+ active budget", SUCCESS),
    ("5", "Snippet\nDelivered", "1\u201310 lines of\nHTML text returned\nmatching page CSS", WARNING),
    ("6", "Injected", "Text inserted into\narticle content at\ndesignated location", WARNING),
    ("7", "Reported", "Impression logged:\ngeo, device, NPI\nspecialty, campaign", ACCENT),
]

for i, (num, title, desc, color) in enumerate(steps):
    x = 0.35 + i * 1.82
    light = PRIMARY_LIGHT if color == PRIMARY else SUCCESS_LIGHT if color == SUCCESS else WARNING_LIGHT if color == WARNING else ACCENT_LIGHT
    add_card(slide, x, 1.85, 1.62, 3.3, light, color)
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x + 0.55), Inches(2.05), Inches(0.45), Inches(0.45))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_text(slide, x + 0.1, 2.6, 1.42, 0.45, title,
             font_size=13, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + 0.05, 3.1, 1.52, 1.6, desc,
             font_size=10, color=BODY_TEXT, alignment=PP_ALIGN.CENTER)
    if i < 6:
        add_text(slide, x + 1.58, 2.9, 0.3, 0.4, "\u2192",
                 font_size=18, color=MUTED, alignment=PP_ALIGN.CENTER)

# NPI branching note
add_card(slide, 0.7, 5.5, 5.5, 1.5, SUCCESS_LIGHT, SUCCESS)
add_text(slide, 0.9, 5.6, 5, 0.3, "If NPI is valid \u2192",
         font_size=13, bold=True, color=SUCCESS)
add_text(slide, 0.9, 5.9, 5, 0.7,
         "Request snippet payload with unit ID and categories. "
         "Engine matches campaign, returns HTML text. "
         "Snippet inherits host page typography.",
         font_size=11, color=BODY_TEXT)

add_card(slide, 6.5, 5.5, 5.8, 1.5, DANGER_LIGHT, DANGER)
add_text(slide, 6.7, 5.6, 5.4, 0.3, "If NPI is missing or invalid \u2192",
         font_size=13, bold=True, color=DANGER)
add_text(slide, 6.7, 5.9, 5.4, 0.7,
         "Unit returns no-op. No snippet is requested, no content is injected, "
         "and no impression is logged. The page renders normally with no visible change.",
         font_size=11, color=BODY_TEXT)

# ════════════════════════════════════════════
# SLIDE 6: SCOPE DEFINITION
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Scope Definition",
             "Clear boundaries to ensure aligned expectations")
page_number(slide, 6)

add_card(slide, 1, 1.85, 5.5, 5.0, SUCCESS_LIGHT, SUCCESS)
add_text(slide, 1.3, 1.95, 5, 0.35, "In Scope \u2014 What Logician Will Build",
         font_size=17, bold=True, color=SUCCESS)
add_line(slide, 1.3, 2.3, 2.5, SUCCESS, 2)
add_bullet_list(slide, 1.3, 2.5, 5, 4.2, [
    "JavaScript tag for page-level deployment",
    "Backend API for NPI verification & campaign matching",
    "HTML snippet delivery system (crawlable, styled)",
    "Category-based campaign routing engine",
    "NPI-based audience targeting integration",
    "Reporting dashboard with filtering and CSV export",
    "Multi-site deployment configuration",
    "Async loading with zero layout shift",
    "HTTPS-only delivery with cache control",
    "Deployment documentation and training",
], font_size=12, color=DARK_TEXT)

add_card(slide, 6.8, 1.85, 5.5, 5.0, CARD_BG, BORDER)
add_text(slide, 7.1, 1.95, 5, 0.35, "Out of Scope \u2014 RxNetwork / Adverge Responsibility",
         font_size=17, bold=True, color=MUTED)
add_line(slide, 7.1, 2.3, 2.5, MUTED, 2)
add_bullet_list(slide, 7.1, 2.5, 5, 4.2, [
    "Medical text snippet content creation",
    "Editorial review and compliance approval",
    "Pharma buyer relationships and sales",
    "NPI identity database (to be provided or integrated)",
    "Campaign budget negotiations with advertisers",
    "Legal / regulatory review of content",
    "Existing ad server (GAM) configuration",
    "Site-level content management decisions",
    "Pharma brand approval workflows",
    "DSP/SSP marketplace setup",
], font_size=12, color=BODY_TEXT)

# ════════════════════════════════════════════
# SLIDE 7: TECHNOLOGY & INFRASTRUCTURE
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Technology & Infrastructure",
             "Purpose-driven stack optimized for speed, reliability, and compliance")
page_number(slide, 7)

stack = [
    ("Client-Side", PRIMARY, PRIMARY_LIGHT, [
        "Vanilla JavaScript \u2014 tag on publisher pages",
        "Zero framework dependencies",
        "Async loading, CSS inheritance",
        "Under 5KB compressed",
    ]),
    ("Backend & API", SUCCESS, SUCCESS_LIGHT, [
        "Python (FastAPI) \u2014 campaign matching",
        "PostgreSQL \u2014 campaigns, impressions",
        "Redis \u2014 caching and fast lookups",
        "RESTful API architecture",
    ]),
    ("Dashboard", WARNING, WARNING_LIGHT, [
        "React / Next.js \u2014 admin panel",
        "Real-time campaign monitoring",
        "Filterable reporting views",
        "CSV/XLSX export capability",
    ]),
    ("Infrastructure", ACCENT, ACCENT_LIGHT, [
        "Cloud-hosted (globally distributed)",
        "HTTPS-only, cache-controlled",
        "Hashed NPI identifiers (no PHI)",
        "Auto-scaling, zero-downtime deploys",
    ]),
]

for i, (title, color, bg, items) in enumerate(stack):
    x = 0.7 + i * 3.1
    add_card(slide, x, 1.85, 2.9, 3.2, bg, color)
    add_text(slide, x + 0.15, 1.95, 2.6, 0.35, title,
             font_size=14, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_line(slide, x + 0.4, 2.35, 2.1, color, 2)
    add_bullet_list(slide, x + 0.1, 2.5, 2.7, 2.4, items,
                    font_size=11, color=BODY_TEXT)

# Security & Compliance row
add_card(slide, 0.7, 5.3, 11.9, 1.8)
add_text(slide, 1, 5.4, 11, 0.35, "Security & Compliance",
         font_size=16, bold=True, color=PRIMARY)

sec_items = [
    ("No PHI Storage", "No protected health information\nis stored or transmitted", PRIMARY),
    ("Hashed NPI", "NPI targeting relies on hashed\nor pseudonymous identifiers", SUCCESS),
    ("HTTPS Only", "All snippet delivery and API\ncommunication over TLS", WARNING),
    ("Cache Control", "Payloads are cache-controlled\nto prevent unintended reuse", ACCENT),
]

for i, (title, desc, color) in enumerate(sec_items):
    x = 0.9 + i * 3.05
    add_text(slide, x, 5.85, 2.8, 0.3, title,
             font_size=12, bold=True, color=color)
    add_text(slide, x, 6.15, 2.8, 0.6, desc,
             font_size=10, color=BODY_TEXT)

# ════════════════════════════════════════════
# SLIDE 8: DELIVERY PHASES & TIMELINE
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Delivery Phases",
             "Structured rollout \u2014 validate early, scale with confidence")
page_number(slide, 8)

# Phase timeline rows
phases = [
    ("Phase 1", "Core System & Single-Site Pilot", "6\u20138 weeks", PRIMARY, PRIMARY_LIGHT,
     ["Backend API (FastAPI + PostgreSQL)",
      "JavaScript tag development & testing",
      "NPI verification integration",
      "Campaign matching engine",
      "Single-site pilot deployment",
      "Basic impression logging"]),
    ("Phase 2", "Multi-Site Rollout & Targeting", "4\u20136 weeks", SUCCESS, SUCCESS_LIGHT,
     ["Deploy across all RxNetwork sites",
      "Full medical category targeting",
      "NPI specialty segmentation",
      "Site-specific configuration",
      "Performance optimization",
      "Frequency capping (if required)"]),
    ("Phase 3", "Dashboard, Reporting & Handover", "4\u20136 weeks", WARNING, WARNING_LIGHT,
     ["Admin dashboard (React/Next.js)",
      "Campaign management interface",
      "Reporting with filters & charts",
      "CSV/XLSX export",
      "Security hardening & audit",
      "Documentation & team training"]),
]

for i, (phase, title, duration, color, bg, items) in enumerate(phases):
    x = 0.7 + i * 4.1
    add_card(slide, x, 1.85, 3.85, 0.7, color, color)
    add_text(slide, x + 0.15, 1.92, 2.5, 0.3, f"{phase}: {title}",
             font_size=13, bold=True, color=WHITE)
    add_text(slide, x + 2.6, 1.92, 1.1, 0.3, duration,
             font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.RIGHT)
    add_card(slide, x, 2.6, 3.85, 3.5, bg, color)
    add_bullet_list(slide, x + 0.15, 2.7, 3.5, 3.2, items,
                    font_size=11, color=DARK_TEXT)

# Timeline summary
add_card(slide, 0.7, 6.35, 11.9, 0.8, CARD_BG, PRIMARY)
add_text(slide, 1, 6.45, 11, 0.25, "Total Timeline: 14\u201320 weeks",
         font_size=16, bold=True, color=PRIMARY, alignment=PP_ALIGN.CENTER)
add_text(slide, 1, 6.75, 11, 0.25,
         "Phase 1 delivers a working system on one site  \u2022  "
         "Each phase builds on the validated output of the previous one",
         font_size=12, color=BODY_TEXT, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# SLIDE 9: INVESTMENT
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Investment",
             "Phase-based pricing \u2014 pay as we deliver, validate as we build")
page_number(slide, 9)

# Phase pricing cards
phase_pricing = [
    ("Phase 1", "Core System & Pilot", "$8,000 \u2013 $12,000",
     "Backend, JS tag, NPI integration,\nsingle-site deployment",
     PRIMARY, PRIMARY_LIGHT),
    ("Phase 2", "Multi-Site Rollout", "$5,000 \u2013 $8,000",
     "All sites, full targeting,\nperformance optimization",
     SUCCESS, SUCCESS_LIGHT),
    ("Phase 3", "Dashboard & Handover", "$5,000 \u2013 $8,000",
     "Admin panel, reporting,\ncompliance, documentation",
     WARNING, WARNING_LIGHT),
]

for i, (phase, title, price, desc, color, bg) in enumerate(phase_pricing):
    x = 0.7 + i * 4.1
    add_card(slide, x, 1.85, 3.85, 3.2, bg, color)
    add_text(slide, x + 0.2, 2.0, 3.45, 0.3, phase,
             font_size=12, bold=True, color=color)
    add_text(slide, x + 0.2, 2.3, 3.45, 0.35, title,
             font_size=15, bold=True, color=DARK_TEXT)
    add_line(slide, x + 0.3, 2.7, 2, color, 2)
    add_text(slide, x + 0.2, 2.85, 3.45, 0.6, price,
             font_size=24, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + 0.2, 3.55, 3.45, 1, desc,
             font_size=12, color=BODY_TEXT, alignment=PP_ALIGN.CENTER)

# Total
add_card(slide, 0.7, 5.3, 11.9, 0.7, PRIMARY_LIGHT, PRIMARY)
add_text(slide, 1, 5.4, 5, 0.4, "Total Estimated Investment",
         font_size=16, bold=True, color=PRIMARY)
add_text(slide, 7, 5.35, 5.3, 0.5, "$18,000 \u2013 $28,000",
         font_size=28, bold=True, color=PRIMARY, alignment=PP_ALIGN.RIGHT)

# Payment structure
add_card(slide, 0.7, 6.2, 11.9, 1.0)
add_text(slide, 1, 6.3, 11, 0.3, "Payment Structure",
         font_size=14, bold=True, color=DARK_TEXT)
add_text(slide, 1, 6.6, 11, 0.4,
         "30% upfront (project kickoff)  \u2022  "
         "30% on Phase 1 delivery  \u2022  "
         "20% on Phase 2 delivery  \u2022  "
         "20% on final delivery & sign-off",
         font_size=12, color=BODY_TEXT)

# ════════════════════════════════════════════
# SLIDE 10: ABOUT LOGICIAN
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "About Logician",
             "Engineering team focused on healthcare ad-tech and data systems")
page_number(slide, 10)

add_card(slide, 1, 1.85, 5.5, 2.5, PRIMARY_LIGHT, PRIMARY)
add_text(slide, 1.3, 2.0, 5, 0.35, "Who We Are",
         font_size=18, bold=True, color=PRIMARY)
add_multiline(slide, 1.3, 2.45, 5, 1.8, [
    "Logician is a focused engineering team that builds data-driven systems for the healthcare and advertising technology industries.",
    "We specialize in programmatic delivery infrastructure, audience targeting systems, and analytics dashboards.",
    "Our approach is phase-based and milestone-driven \u2014 you see working results early and pay as we deliver.",
], font_size=12, color=BODY_TEXT)

add_card(slide, 6.8, 1.85, 5.5, 2.5, SUCCESS_LIGHT, SUCCESS)
add_text(slide, 7.1, 2.0, 5, 0.35, "Relevant Capabilities",
         font_size=18, bold=True, color=SUCCESS)
add_bullet_list(slide, 7.1, 2.45, 5, 1.8, [
    "Programmatic ad delivery systems",
    "NPI and healthcare audience targeting",
    "Real-time data dashboards and reporting",
    "Scalable API development and cloud infrastructure",
    "Multi-site deployment and configuration management",
], font_size=12, color=DARK_TEXT)

add_card(slide, 1, 4.65, 11.3, 2.5)
add_text(slide, 1.3, 4.8, 10, 0.35, "How We Work",
         font_size=18, bold=True, color=PRIMARY)

work_items = [
    ("Phase-Based Delivery",
     "Every phase produces a working, testable system. You validate before we expand.",
     PRIMARY),
    ("Weekly Check-ins",
     "Short weekly calls to demo progress, review priorities, and resolve blockers.",
     SUCCESS),
    ("Milestone Payments",
     "You pay tied to deliverables \u2014 not time elapsed. If we don\u2019t deliver, you don\u2019t pay.",
     WARNING),
    ("Full Documentation",
     "Every system we build includes deployment guides, API docs, and training materials.",
     ACCENT),
]

for i, (title, desc, color) in enumerate(work_items):
    x = 1.3 + i * 2.8
    add_text(slide, x, 5.25, 2.5, 0.3, title,
             font_size=12, bold=True, color=color)
    add_text(slide, x, 5.55, 2.5, 1.2, desc,
             font_size=11, color=BODY_TEXT)

# ════════════════════════════════════════════
# SLIDE 11: NEXT STEPS
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Next Steps",
             "What we need to finalize scope and begin delivery")
page_number(slide, 11)

add_card(slide, 1, 1.85, 5.5, 4.5, PRIMARY_LIGHT, PRIMARY)
add_text(slide, 1.3, 1.95, 5, 0.35, "Questions We Need Answered",
         font_size=17, bold=True, color=PRIMARY)
add_line(slide, 1.3, 2.3, 2, PRIMARY, 2)
add_bullet_list(slide, 1.3, 2.5, 5, 3.5, [
    "How is NPI identity currently resolved on your sites?",
    "Is there an existing login system for HCPs?",
    "What ad-serving infrastructure does Adverge use today?",
    "How many sites are in scope for Phase 1?",
    "Approximate monthly traffic volume?",
    "Who creates and approves the text snippet content?",
    "What is your target launch date?",
    "Is there a defined viewability standard for text units?",
], font_size=12, color=DARK_TEXT)

add_card(slide, 6.8, 1.85, 5.5, 4.5, SUCCESS_LIGHT, SUCCESS)
add_text(slide, 7.1, 1.95, 5, 0.35, "Once Confirmed \u2014 We Begin",
         font_size=17, bold=True, color=SUCCESS)
add_line(slide, 7.1, 2.3, 2, SUCCESS, 2)
add_bullet_list(slide, 7.1, 2.5, 5, 3.5, [
    "Finalize scope based on your answers",
    "Sign agreement and confirm payment schedule",
    "Kick off Phase 1: database design and API planning",
    "Set up shared communication channel",
    "Schedule weekly progress demos",
    "Deliver working single-site pilot within 6\u20138 weeks",
    "Review Phase 1 results before approving Phase 2",
    "Full system live within 14\u201320 weeks",
], font_size=12, color=DARK_TEXT)

add_card(slide, 1, 6.65, 11.3, 0.6, CARD_BG, PRIMARY)
add_text(slide, 1.3, 6.7, 10, 0.4,
         "We\u2019re ready to begin. Let\u2019s schedule a discovery call to align on scope and timeline.",
         font_size=15, bold=True, color=PRIMARY, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════
output_path = r"c:\Users\akash\work\barron\RxNetwork_Proposal_Logician.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
